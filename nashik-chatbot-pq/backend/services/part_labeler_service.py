"""
PartLabeler Service
Handles business logic for CAD image labeling and warranty lookup
"""

import logging
import os
import pandas as pd
from datetime import datetime
from typing import List, Dict, Any, Optional
from app.connectors.state_db_connector import StateDBConnector
from app.queries.part_labeler_queries import PartLabelerQueries
from sqlalchemy import text

logger = logging.getLogger(__name__)

QUARTER_MONTH_MAP = {1: ('Jan', 'Mar'), 2: ('Jan', 'Mar'), 3: ('Jan', 'Mar'),
                     4: ('Apr', 'Jun'), 5: ('Apr', 'Jun'), 6: ('Apr', 'Jun'),
                     7: ('Jul', 'Sep'), 8: ('Jul', 'Sep'), 9: ('Jul', 'Sep'),
                     10: ('Oct', 'Dec'), 11: ('Oct', 'Dec'), 12: ('Oct', 'Dec')}

def derive_mfg_month(date_val: str) -> str:
    """Convert a date string (YYYY-MM-DD or YYYY-MM-DD HH:MM:SS) to 'Mon-YY' format."""
    if not date_val or str(date_val).strip().lower() in ('nan', '', 'none', 'nat'):
        return ''
    try:
        date_str = str(date_val).split(' ')[0].strip()  # take YYYY-MM-DD part
        dt = datetime.strptime(date_str, '%Y-%m-%d')
        return dt.strftime('%b-%y')  # e.g. "Jan-26"
    except Exception:
        return ''

def format_mfg_month_warranty(val: str) -> str:
    """Convert any date/month value to 'Mon-YYYY' format required by warranty SQL queries."""
    if not val or str(val).strip().lower() in ('nan', '', 'none', 'nat'):
        return ''
    val = str(val).strip()
    # Already in Mon-YYYY or Mon-YY format — normalize to Mon-YYYY
    for fmt in ('%b-%Y', '%b-%y'):
        try:
            return datetime.strptime(val, fmt).strftime('%b-%Y')
        except ValueError:
            pass
    # Date string like "2025-01-27 00:00:00" or "2025-01-27"
    date_str = val.split(' ')[0].strip()
    for fmt in ('%Y-%m-%d', '%d-%m-%Y', '%m/%d/%Y', '%d/%m/%Y'):
        try:
            return datetime.strptime(date_str, fmt).strftime('%b-%Y')
        except ValueError:
            pass
    return val

def derive_mfg_quarter(date_val: str) -> str:
    """Convert a date string to quarter label e.g. 'Jan26-Mar26'."""
    if not date_val or str(date_val).strip().lower() in ('nan', '', 'none', 'nat'):
        return ''
    try:
        date_str = str(date_val).split(' ')[0].strip()
        dt = datetime.strptime(date_str, '%Y-%m-%d')
        start_mon, end_mon = QUARTER_MONTH_MAP[dt.month]
        yy = dt.strftime('%y')
        return f"{start_mon}{yy}-{end_mon}{yy}"
    except Exception:
        return ''

def safe_str(val) -> str:
    """Convert value to string, treating NaN/None as empty."""
    if val is None:
        return ''
    s = str(val).strip()
    return '' if s.lower() in ('nan', 'none', 'nat') else s

def normalize_month_label(label: str) -> str:
    """Normalize any month label to 'Mon-YY' (e.g. 'Jan-2024' -> 'Jan-24')."""
    if not label:
        return label
    for fmt in ('%b-%Y', '%b-%y'):
        try:
            return datetime.strptime(label, fmt).strftime('%b-%y')
        except ValueError:
            pass
    return label

def generate_month_sequence(min_date, max_date) -> list:
    """Build complete list of 'Mon-YY' labels from min_date to max_date (inclusive)."""
    sequence = []
    curr = datetime(min_date.year, min_date.month, 1)
    end = datetime(max_date.year, max_date.month, 1)
    while curr <= end:
        sequence.append(curr.strftime('%b-%y'))
        m = curr.month + 1
        y = curr.year
        if m > 12:
            m = 1
            y += 1
        curr = datetime(y, m, 1)
    return sequence

class PartLabelerService:
    """Manages CAD labeling data and warranty information"""

    def __init__(self):
        self.db = StateDBConnector()

    def extract_excel_headers(self, file_path: str) -> List[str]:
        """Read the first row of Excel/CSV and return headers"""
        try:
            if file_path.endswith('.csv'):
                df = pd.read_csv(file_path, nrows=0)
            else:
                df = pd.read_excel(file_path, nrows=0)
            return df.columns.tolist()
        except Exception as e:
            logger.error(f"Error extracting headers: {e}")
            raise Exception(f"Failed to read file headers: {str(e)}")

    def process_mapped_warranty_data(self, file_path: str, mapping: Dict[str, str], user_id: int) -> int:
        """
        Process the file using the provided column mapping.
        mapping format: { 'internal_db_col': 'user_excel_col' }
        """
        try:
            if file_path.endswith('.csv'):
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path)

            # Clean and prepare records
            records = []
            
            # Clear existing data for this user before loading new one
            self.db.execute_update("DELETE FROM raw_warranty_data WHERE user_id = :user_id", {"user_id": user_id})

            for _, row in df.iterrows():
                record = {}
                # Initialize all 44 columns with empty string to ensure SQL parameters match
                all_cols = [
                    'region', 'zone', 'area_office', 'plant', 'plant_desc', 'commodity', 
                    'group_code', 'group_code_desc', 'complaint_code', 'complaint_code_desc',
                    'base_model', 'model_code', 'model_family', 'claim_type', 'sap_claim_no', 
                    'claim_desc', 'ac_non_ac', 'variant', 'drive_type', 'service_type', 
                    'billing_dealer', 'billing_dealer_name', 'serial_no', 'claim_date', 
                    'failure_kms', 'km_hr_group', 'dealer_verbatim', 'part', 'vender', 
                    'material_description', 'causal_flag', 'jdp_city', 'fisyr_qrt', 
                    'engine_number', 'manufac_yr_mon', 'failure_date', 'mis_bucket', 
                    'walk_home', 'dealer_code', 'claim_dealer_name', 'ro_number', 
                    'no_of_incidents', 'new_manufacturing_quater', 'vendor_manuf'
                ]
                for col in all_cols:
                    record[col] = ''
                
                record['user_id'] = user_id

                # Apply mapping
                for db_col, user_col in mapping.items():
                    if user_col in row:
                        val = str(row[user_col])
                        if val.lower() == 'nan': val = ''
                        record[db_col] = val.strip()
                
                # Normalize manufac_yr_mon to 'Mon-YYYY' format (e.g. "Jan-2025")
                # Excel date columns come in as datetime strings like "2025-01-27 00:00:00"
                # which break TO_DATE(manufac_yr_mon, 'Mon-YYYY') in SQL queries
                if record.get('manufac_yr_mon'):
                    record['manufac_yr_mon'] = format_mfg_month_warranty(record['manufac_yr_mon'])

                # Special handling: if failure_date is not mapped but claim_date is,
                # use claim_date as failure_date fallback for the trend logic
                if not record.get('failure_date') and record.get('claim_date'):
                    record['failure_date'] = record['claim_date']

                records.append(record)

                if len(records) >= 500:
                    self._insert_batch(records)
                    records = []

            if records:
                self._insert_batch(records)

            return len(df)
        except Exception as e:
            logger.error(f"Error processing mapped data: {e}")
            raise Exception(f"Failed to process data: {str(e)}")

    def _insert_batch(self, records):
        with self.db.get_session() as session:
            session.execute(text(PartLabelerQueries.INSERT_RAW_WARRANTY), records)

    def get_warranty_data(self, user_id: int, part_name: Optional[str] = None, month: Optional[list[str]] = None, base_model: Optional[list[str]] = None, mis_bucket: Optional[list[str]] = None, mfg_qtr: Optional[list[str]] = None) -> Any:
        """
        Lookup warranty data from raw_warranty_data table
        """
        try:
            params = {
                "base_model": base_model if base_model and "All" not in base_model else None,
                "mis_bucket": mis_bucket if mis_bucket and "All" not in mis_bucket else None,
                "mfg_qtr": mfg_qtr if mfg_qtr and "All" not in mfg_qtr else None,
                "user_id": user_id
            }
            
            if part_name:
                normalized_search = f"%{part_name.lower().replace(' ', '')}%"
                params["search_term"] = normalized_search
                rows = self.db.execute_query(
                    PartLabelerQueries.SEARCH_WARRANTY_DATA,
                    params
                )
            else:
                rows = self.db.execute_query(PartLabelerQueries.GET_WARRANTY_DATA, params)

            # Map rows to dicts
            data = [
                {
                    "partName": r[0],
                    "month": r[1],
                    "failureCount": r[2],
                    "description": r[3]
                }
                for r in rows
            ]

            # Filter by month in memory if needed
            if month and "All" not in month:
                # Use exact match as provided (e.g., "Dec-2024")
                data = [item for item in data if item['month'] in month]

            return data
        except Exception as e:
            logger.error(f"Error fetching warranty data from DB: {e}")
            return {"error": str(e)}

    def get_filter_options(self, user_id: int) -> Dict[str, List[str]]:
        """Fetch unique models, MIS buckets, Mfg Quarters, and Months for filters"""
        try:
            params = {"user_id": user_id}
            models = [r[0] for r in self.db.execute_query(PartLabelerQueries.GET_UNIQUE_MODELS, params)]
            mis_buckets = [r[0] for r in self.db.execute_query(PartLabelerQueries.GET_UNIQUE_MIS_BUCKETS, params)]
            mfg_quarters = [r[0] for r in self.db.execute_query(PartLabelerQueries.GET_UNIQUE_MFG_QUARTERS, params)]
            mfg_months = [r[0] for r in self.db.execute_query(PartLabelerQueries.GET_UNIQUE_MFG_MONTHS, params)]
            return {
                "models": models, 
                "mis_buckets": mis_buckets, 
                "mfg_quarters": mfg_quarters,
                "mfg_months": mfg_months
            }
        except Exception as e:
            logger.error(f"Error fetching filter options: {e}")
            return {"models": [], "mis_buckets": [], "mfg_quarters": [], "mfg_months": []}

    def get_detailed_warranty_csv(self, user_id: int, part_name: str, month: Optional[list[str]] = None, base_model: Optional[list[str]] = None, mis_bucket: Optional[list[str]] = None, mfg_qtr: Optional[list[str]] = None) -> str:
        """Fetch all columns for a part and return as CSV string"""
        import csv
        import io
        try:
            normalized_search = f"%{part_name.lower().replace(' ', '')}%"
            params = {
                "search_term": normalized_search,
                "base_model": base_model if base_model and "All" not in base_model else None,
                "mis_bucket": mis_bucket if mis_bucket and "All" not in mis_bucket else None,
                "mfg_qtr": mfg_qtr if mfg_qtr and "All" not in mfg_qtr else None,
                "month_val": month if month and "All" not in month else None,
                "user_id": user_id
            }
            
            # Use the new method to get both headers and rows
            headers, rows = self.db.execute_query_with_headers(
                PartLabelerQueries.GET_ALL_WARRANTY_FOR_PART,
                params
            )
            
            if not rows:
                return "No data found"

            output = io.StringIO()
            writer = csv.writer(output)
            
            # Write column names first
            writer.writerow(headers)
            # Write data rows
            writer.writerows(rows)
            return output.getvalue()
        except Exception as e:
            logger.error(f"Error generating CSV: {e}")
            return f"Error generating CSV: {str(e)}"

    def get_detailed_warranty_pdf(self, user_id: int, part_name: str, month: Optional[list[str]] = None, base_model: Optional[list[str]] = None, mis_bucket: Optional[list[str]] = None, mfg_qtr: Optional[list[str]] = None) -> bytes:
        """Fetch all columns for a part and return as a PDF (bytes)"""
        import io
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import mm
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer

        try:
            normalized_search = f"%{part_name.lower().replace(' ', '')}%"
            params = {
                "search_term": normalized_search,
                "base_model": base_model if base_model and "All" not in base_model else None,
                "mis_bucket": mis_bucket if mis_bucket and "All" not in mis_bucket else None,
                "mfg_qtr": mfg_qtr if mfg_qtr and "All" not in mfg_qtr else None,
                "month_val": month if month and "All" not in month else None,
                "user_id": user_id
            }

            headers, rows = self.db.execute_query_with_headers(
                PartLabelerQueries.GET_ALL_WARRANTY_FOR_PART,
                params
            )

            output = io.BytesIO()
            doc = SimpleDocTemplate(
                output,
                pagesize=landscape(A4),
                leftMargin=12 * mm, rightMargin=12 * mm,
                topMargin=12 * mm, bottomMargin=12 * mm
            )
            styles = getSampleStyleSheet()
            elements = [
                Paragraph(f"Warranty Data - {part_name}", styles["Title"]),
                Spacer(1, 6)
            ]

            if not rows:
                elements.append(Paragraph("No data found", styles["Normal"]))
            else:
                # Keep cells readable: wrap header/body text in paragraphs
                cell_style = styles["BodyText"]
                cell_style.fontSize = 5.5
                cell_style.leading = 6.5
                table_data = [[Paragraph(str(h), cell_style) for h in headers]]
                for row in rows:
                    table_data.append([Paragraph("" if v is None else str(v), cell_style) for v in row])

                # reportlab does NOT auto-shrink columns to fit the page like PPT does -
                # without explicit colWidths a wide table (many columns) raises a
                # LayoutError ("Flowable too large on page") instead of rendering.
                # Force every column into an equal share of the available page width.
                available_width = landscape(A4)[0] - doc.leftMargin - doc.rightMargin
                col_width = available_width / len(headers)
                table = Table(table_data, colWidths=[col_width] * len(headers), repeatRows=1)
                table.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#DC0028")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("FONTSIZE", (0, 0), (-1, 0), 5.5),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#e2e8f0")),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f8fafc")]),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 2),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 2),
                    ("TOPPADDING", (0, 0), (-1, -1), 2),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 2),
                ]))
                elements.append(table)

            doc.build(elements)
            return output.getvalue()
        except Exception as e:
            logger.error(f"Error generating PDF: {e}")
            raise

    def get_detailed_warranty_pptx(self, user_id: int, part_name: str, month: Optional[list[str]] = None, base_model: Optional[list[str]] = None, mis_bucket: Optional[list[str]] = None, mfg_qtr: Optional[list[str]] = None, rows_per_slide: int = 12) -> bytes:
        """Fetch all columns for a part and return as a PPTX (bytes), one table per slide"""
        import io
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        HEADER_RED = RGBColor(0xDC, 0x00, 0x28)
        STRIPE_GRAY = RGBColor(0xF8, 0xFA, 0xFC)

        try:
            normalized_search = f"%{part_name.lower().replace(' ', '')}%"
            params = {
                "search_term": normalized_search,
                "base_model": base_model if base_model and "All" not in base_model else None,
                "mis_bucket": mis_bucket if mis_bucket and "All" not in mis_bucket else None,
                "mfg_qtr": mfg_qtr if mfg_qtr and "All" not in mfg_qtr else None,
                "month_val": month if month and "All" not in month else None,
                "user_id": user_id
            }

            headers, rows = self.db.execute_query_with_headers(
                PartLabelerQueries.GET_ALL_WARRANTY_FOR_PART,
                params
            )

            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            blank_layout = prs.slide_layouts[6]

            # Title slide
            title_slide = prs.slides.add_slide(blank_layout)
            title_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.33), Inches(1.5))
            tf = title_box.text_frame
            tf.text = f"Warranty Data - {part_name}"
            tf.paragraphs[0].font.size = Pt(32)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            sub = tf.add_paragraph()
            sub.text = f"{len(rows)} record(s)"
            sub.font.size = Pt(16)
            sub.alignment = PP_ALIGN.CENTER

            if not rows:
                empty_slide = prs.slides.add_slide(blank_layout)
                box = empty_slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.33), Inches(1))
                box.text_frame.text = "No data found"
                box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            else:
                n_cols = len(headers)
                for start in range(0, len(rows), rows_per_slide):
                    chunk = rows[start:start + rows_per_slide]
                    slide = prs.slides.add_slide(blank_layout)
                    n_rows = len(chunk) + 1  # + header
                    table_shape = slide.shapes.add_table(
                        n_rows, n_cols,
                        Inches(0.2), Inches(0.3),
                        Inches(12.93), Inches(6.9)
                    )
                    table = table_shape.table

                    for col_idx, header in enumerate(headers):
                        cell = table.cell(0, col_idx)
                        cell.text = str(header)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = HEADER_RED
                        for p in cell.text_frame.paragraphs:
                            p.font.size = Pt(8)
                            p.font.bold = True
                            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

                    for r, row in enumerate(chunk, start=1):
                        for c, value in enumerate(row):
                            cell = table.cell(r, c)
                            cell.text = "" if value is None else str(value)
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = STRIPE_GRAY if r % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
                            for p in cell.text_frame.paragraphs:
                                p.font.size = Pt(7)

            output = io.BytesIO()
            prs.save(output)
            return output.getvalue()
        except Exception as e:
            logger.error(f"Error generating PPTX: {e}")
            raise

    def upload_image(self, filename: str, user_id: int, display_name: Optional[str] = None) -> Dict[str, Any]:
        """Save image record to database"""
        image_id = self.db.execute_insert(
            PartLabelerQueries.INSERT_IMAGE,
            {"filename": filename, "display_name": display_name or filename, "user_id": user_id}
        )
        return {"id": image_id, "filename": filename, "display_name": display_name or filename}

    def get_all_images(self, user_id: int) -> List[Dict[str, Any]]:
        """Get all uploaded images"""
        rows = self.db.execute_query(PartLabelerQueries.GET_ALL_IMAGES, {"user_id": user_id})
        return [{"id": r[0], "filename": r[1], "created_at": r[2], "display_name": r[3] or r[1]} for r in rows]

    def delete_image(self, image_id: int, user_id: int) -> bool:
        """Delete an image and its labels"""
        self.db.execute_update(PartLabelerQueries.DELETE_IMAGE, {"id": image_id, "user_id": user_id})
        return True

    def save_label(self, label_data: Dict[str, Any], user_id: int) -> bool:
        """Save a new label marker"""
        self.db.execute_insert(
            PartLabelerQueries.INSERT_LABEL,
            {
                "image_id": label_data["imageId"],
                "part_name": label_data["partName"],
                "description": label_data.get("description", ""),
                "part_number": label_data.get("partNumber", ""),
                "failure_count": label_data.get("failureCount", 0),
                "report_month": label_data.get("reportMonth", "All"),
                "x_coord": label_data["x"],
                "y_coord": label_data["y"],
                "user_id": user_id
            }
        )
        return True

    def delete_label(self, label_id: int, user_id: int) -> bool:
        """Delete a label marker"""
        self.db.execute_update(PartLabelerQueries.DELETE_LABEL, {"id": label_id, "user_id": user_id})
        return True

    def update_label_name(self, label_id: int, part_name: str, user_id: int) -> bool:
        """Update part name for a label"""
        self.db.execute_update(
            PartLabelerQueries.UPDATE_LABEL_NAME,
            {"id": label_id, "part_name": part_name, "user_id": user_id}
        )
        return True

    def get_labels_for_image(self, image_id: int, user_id: int) -> List[Dict[str, Any]]:
        """Get all labels associated with an image"""
        rows = self.db.execute_query(
            PartLabelerQueries.GET_LABELS_FOR_IMAGE,
            {"image_id": image_id, "user_id": user_id}
        )
        return [
            {
                "id": r[0],
                "imageId": r[1],
                "partName": r[2],
                "description": r[3],
                "partNumber": r[4],
                "failureCount": r[5],
                "reportMonth": r[6],
                "x": r[7],
                "y": r[8]
            }
            for r in rows
        ]

    # =====================================================
    # DATA SOURCE DISPATCHER METHODS
    # =====================================================

    def get_data_status(self, user_id: int) -> dict:
        """Return which of the data sources have rows uploaded for this user."""
        sources = {
            "warranty": "raw_warranty_data",
            "rpt":      "raw_rpt_data",
            "gnovac":   "raw_gnovac_data",
            "rfi":      "raw_rfi_data",
            "esqa":     "raw_esqa_data",
            "ev":       "raw_ev_data",
        }
        status = {}
        for key, table in sources.items():
            rows = self.db.execute_query(
                f"SELECT COUNT(*) FROM {table} WHERE user_id = :user_id",
                {"user_id": user_id},
            )
            count = rows[0][0] if rows else 0
            status[key] = {"uploaded": count > 0, "row_count": count}
        return status

    def process_data_for_source(self, file_path: str, mapping: Dict[str, str], user_id: int, data_source: str) -> int:
        """Dispatch file processing to the correct ingestion method based on data_source."""
        if data_source == 'rpt':
            return self.process_mapped_rpt_data(file_path, mapping, user_id)
        elif data_source == 'gnovac':
            return self.process_mapped_gnovac_data(file_path, mapping, user_id)
        elif data_source == 'rfi':
            return self.process_mapped_rfi_data(file_path, mapping, user_id)
        elif data_source == 'esqa':
            return self.process_mapped_esqa_data(file_path, mapping, user_id)
        elif data_source == 'ev':
            return self.process_mapped_ev_data(file_path, mapping, user_id)
        else:
            return self.process_mapped_warranty_data(file_path, mapping, user_id)

    def get_filter_options_for_source(self, user_id: int, data_source: str) -> Dict[str, List[str]]:
        """Dispatch filter options to the correct method based on data_source."""
        if data_source == 'rpt':
            return self.get_rpt_filter_options(user_id)
        elif data_source == 'gnovac':
            return self.get_gnovac_filter_options(user_id)
        elif data_source == 'rfi':
            return self.get_rfi_filter_options(user_id)
        elif data_source == 'esqa':
            return self.get_esqa_filter_options(user_id)
        elif data_source == 'ev':
            return self.get_ev_filter_options(user_id)
        else:
            return self.get_filter_options(user_id)

    def get_source_data(self, user_id: int, part_name: Optional[str], month, base_model, mis_bucket, mfg_qtr, data_source: str, buyoff_stage=None, online_offline=None, defect_type=None, battery_motor=None) -> Any:
        """Dispatch data lookup to the correct method based on data_source."""
        if data_source == 'rpt':
            return self.get_rpt_data(user_id, part_name, month, base_model, mis_bucket, mfg_qtr, buyoff_stage, online_offline)
        elif data_source == 'gnovac':
            return self.get_gnovac_data(user_id, part_name, month, base_model, mis_bucket, mfg_qtr)
        elif data_source == 'rfi':
            return self.get_rfi_data(user_id, part_name, month, base_model, mis_bucket, mfg_qtr, defect_type)
        elif data_source == 'esqa':
            return self.get_esqa_data(user_id, part_name, month, base_model, mis_bucket, mfg_qtr)
        elif data_source == 'ev':
            return self.get_ev_data(user_id, part_name, month, base_model, mis_bucket, mfg_qtr, battery_motor)
        else:
            return self.get_warranty_data(user_id, part_name, month, base_model, mis_bucket, mfg_qtr)

    def get_dashboard_data_for_source(self, user_id: int, part_name, month, base_model, mis_bucket, mfg_qtr, data_source: str, buyoff_stage=None, online_offline=None, defect_type=None, battery_motor=None) -> Dict[str, Any]:
        """Dispatch dashboard data to the correct method based on data_source."""
        if data_source == 'rpt':
            return self.get_rpt_dashboard_data(user_id, part_name, month, base_model, mis_bucket, mfg_qtr, buyoff_stage, online_offline)
        elif data_source == 'gnovac':
            return self.get_gnovac_dashboard_data(user_id, part_name, month, base_model, mis_bucket, mfg_qtr)
        elif data_source == 'rfi':
            return self.get_rfi_dashboard_data(user_id, part_name, month, base_model, mis_bucket, mfg_qtr, defect_type)
        elif data_source == 'esqa':
            return self.get_esqa_dashboard_data(user_id, part_name, month, base_model, mis_bucket, mfg_qtr)
        elif data_source == 'ev':
            return self.get_ev_dashboard_data(user_id, part_name, month, base_model, mis_bucket, mfg_qtr, battery_motor)
        else:
            return self.get_dashboard_data(user_id, part_name, month, base_model, mis_bucket, mfg_qtr)

    # =====================================================
    # OFFLINE RPT METHODS
    # =====================================================

    def process_mapped_rpt_data(self, file_path: str, mapping: Dict[str, str], user_id: int) -> int:
        try:
            df = pd.read_csv(file_path) if file_path.endswith('.csv') else pd.read_excel(file_path)
            all_cols = [
                'date_col', 'mfg_month', 'mfg_quarter', 'shift', 'body_sr_no', 'vin_number',
                'buyoff_stage', 'model', 'platform_group', 'stage_name', 'part', 'defect',
                'part_defect', 'attribute_name', 'custom_attribution', 'offline_val', 'online_val',
                'rework_status', 'location_name', 'defect_status', 'as_is_ok', 'shop_name',
                'model_description', 'model_code', 'severity_name', 'domestic_export', 'defect_category'
            ]
            self.db.execute_update("DELETE FROM raw_rpt_data WHERE user_id = :user_id", {"user_id": user_id})
            records = []
            for _, row in df.iterrows():
                record = {col: '' for col in all_cols}
                record['user_id'] = user_id
                for db_col, user_col in mapping.items():
                    if user_col in row:
                        record[db_col] = safe_str(row[user_col])
                # Derive mfg_month and mfg_quarter from date_col
                record['mfg_month'] = derive_mfg_month(record['date_col'])
                record['mfg_quarter'] = derive_mfg_quarter(record['date_col'])
                records.append(record)
                if len(records) >= 500:
                    self._insert_batch_rpt(records)
                    records = []
            if records:
                self._insert_batch_rpt(records)
            return len(df)
        except Exception as e:
            logger.error(f"Error processing RPT data: {e}")
            raise

    def _insert_batch_rpt(self, records):
        with self.db.get_session() as session:
            session.execute(text(PartLabelerQueries.INSERT_RAW_RPT), records)

    def get_rpt_filter_options(self, user_id: int) -> Dict[str, List[str]]:
        params = {"user_id": user_id}
        return {
            "models": [r[0] for r in self.db.execute_query(PartLabelerQueries.RPT_GET_UNIQUE_MODELS, params)],
            "mis_buckets": [r[0] for r in self.db.execute_query(PartLabelerQueries.RPT_GET_UNIQUE_MIS, params)],
            "mfg_quarters": [r[0] for r in self.db.execute_query(PartLabelerQueries.RPT_GET_UNIQUE_MFG_QUARTERS, params)],
            "mfg_months": [r[0] for r in self.db.execute_query(PartLabelerQueries.RPT_GET_UNIQUE_MFG_MONTHS, params)],
            "buyoff_stages": [r[0] for r in self.db.execute_query(PartLabelerQueries.RPT_GET_UNIQUE_BUYOFF_STAGES, params)],
            "online_offline_options": ["Online", "Offline"],
        }

    def get_rpt_data(self, user_id: int, part_name=None, month=None, base_model=None, mis_bucket=None, mfg_qtr=None, buyoff_stage=None, online_offline=None) -> Any:
        try:
            params = {
                "base_model": base_model if base_model and "All" not in base_model else None,
                "mis_bucket": mis_bucket if mis_bucket and "All" not in mis_bucket else None,
                "mfg_qtr": mfg_qtr if mfg_qtr and "All" not in mfg_qtr else None,
                "buyoff_stage": buyoff_stage if buyoff_stage and "All" not in buyoff_stage else None,
                "online_offline": online_offline if online_offline and "All" not in online_offline else None,
                "user_id": user_id,
            }
            if part_name:
                params["search_term"] = f"%{part_name.lower().replace(' ', '')}%"
                rows = self.db.execute_query(PartLabelerQueries.RPT_SEARCH_DATA, params)
            else:
                rows = []
            data = [{"partName": r[0], "month": r[1], "failureCount": r[2], "description": r[3]} for r in rows]
            if month and "All" not in month:
                data = [d for d in data if d['month'] in month]
            return data
        except Exception as e:
            logger.error(f"RPT data lookup error: {e}")
            return {"error": str(e)}

    def get_rpt_dashboard_data(self, user_id: int, part_name=None, month=None, base_model=None, mis_bucket=None, mfg_qtr=None, buyoff_stage=None, online_offline=None) -> Dict[str, Any]:
        result = {"mfgMonth": [], "reportingMonth": [], "kms": [], "region": []}
        search_terms = [f"%{p.lower().replace(' ', '')}%" for p in part_name if p] if part_name else None
        params = {
            "user_id": user_id,
            "month_val": month if month and "All" not in month else None,
            "base_model": base_model if base_model and "All" not in base_model else None,
            "mis_bucket": mis_bucket if mis_bucket and "All" not in mis_bucket else None,
            "mfg_qtr": mfg_qtr if mfg_qtr and "All" not in mfg_qtr else None,
            "buyoff_stage": buyoff_stage if buyoff_stage and "All" not in buyoff_stage else None,
            "online_offline": online_offline if online_offline and "All" not in online_offline else None,
            "search_terms": search_terms,
        }
        try:
            range_row = self.db.execute_query(PartLabelerQueries.RPT_GET_MFG_DATE_RANGE, {"user_id": user_id})
            if range_row and range_row[0][0] and range_row[0][1]:
                min_date, max_date = range_row[0]
                sequence = generate_month_sequence(min_date, max_date)
                db_data = {}
                if search_terms:
                    db_data = {r[0]: r[1] for r in self.db.execute_query(PartLabelerQueries.RPT_GET_DASHBOARD_MFG_MONTH, params)}
                result["mfgMonth"] = [{"label": m, "value": db_data.get(m, 0)} for m in sequence]
        except Exception as e:
            logger.error(f"RPT mfgMonth error: {e}")
        try:
            result["reportingMonth"] = [{"label": r[0], "value": r[1]} for r in self.db.execute_query(PartLabelerQueries.RPT_GET_DASHBOARD_REPORTING_MONTH, params)]
        except Exception as e:
            logger.error(f"RPT reportingMonth error: {e}")
        try:
            result["kms"] = [{"label": r[0], "value": r[1]} for r in self.db.execute_query(PartLabelerQueries.RPT_GET_DASHBOARD_SHIFT, params)]
        except Exception as e:
            logger.error(f"RPT shift error: {e}")
        try:
            result["region"] = [{"label": r[0], "value": r[1]} for r in self.db.execute_query(PartLabelerQueries.RPT_GET_DASHBOARD_LOCATION, params)]
        except Exception as e:
            logger.error(f"RPT location error: {e}")
        return result

    # =====================================================
    # GNOVAC METHODS
    # =====================================================

    def process_mapped_gnovac_data(self, file_path: str, mapping: Dict[str, str], user_id: int) -> int:
        try:
            df = pd.read_csv(file_path) if file_path.endswith('.csv') else pd.read_excel(file_path)
            all_cols = [
                'audit_date', 'mfg_month', 'mfg_quarter', 'vin_no', 'plant_name', 'model_code',
                'variant_name', 'fuel_type', 'build_phase_name', 'body_no', 'part_name', 'defect_name',
                'location_name', 'concern_type_name', 'pointer', 'attribution', 'four_m',
                'four_m_analysis_name', 'root_cause', 'ica', 'pca', 'responsibility', 'target_date',
                'status', 'frequency', 'new_and_repeat', 'remark'
            ]
            self.db.execute_update("DELETE FROM raw_gnovac_data WHERE user_id = :user_id", {"user_id": user_id})
            records = []
            for _, row in df.iterrows():
                record = {col: '' for col in all_cols}
                record['user_id'] = user_id
                for db_col, user_col in mapping.items():
                    if user_col in row:
                        record[db_col] = safe_str(row[user_col])
                record['mfg_month'] = derive_mfg_month(record['audit_date'])
                record['mfg_quarter'] = derive_mfg_quarter(record['audit_date'])
                records.append(record)
                if len(records) >= 500:
                    self._insert_batch_gnovac(records)
                    records = []
            if records:
                self._insert_batch_gnovac(records)
            return len(df)
        except Exception as e:
            logger.error(f"Error processing GNOVAC data: {e}")
            raise

    def _insert_batch_gnovac(self, records):
        with self.db.get_session() as session:
            session.execute(text(PartLabelerQueries.INSERT_RAW_GNOVAC), records)

    def get_gnovac_filter_options(self, user_id: int) -> Dict[str, List[str]]:
        params = {"user_id": user_id}
        return {
            "models": [r[0] for r in self.db.execute_query(PartLabelerQueries.GNOVAC_GET_UNIQUE_MODELS, params)],
            "mis_buckets": [r[0] for r in self.db.execute_query(PartLabelerQueries.GNOVAC_GET_UNIQUE_MIS, params)],
            "mfg_quarters": [r[0] for r in self.db.execute_query(PartLabelerQueries.GNOVAC_GET_UNIQUE_MFG_QUARTERS, params)],
            "mfg_months": [r[0] for r in self.db.execute_query(PartLabelerQueries.GNOVAC_GET_UNIQUE_MFG_MONTHS, params)],
        }

    def get_gnovac_data(self, user_id: int, part_name=None, month=None, base_model=None, mis_bucket=None, mfg_qtr=None) -> Any:
        try:
            params = {
                "base_model": base_model if base_model and "All" not in base_model else None,
                "mis_bucket": mis_bucket if mis_bucket and "All" not in mis_bucket else None,
                "mfg_qtr": mfg_qtr if mfg_qtr and "All" not in mfg_qtr else None,
                "user_id": user_id,
            }
            if part_name:
                params["search_term"] = f"%{part_name.lower().replace(' ', '')}%"
                rows = self.db.execute_query(PartLabelerQueries.GNOVAC_SEARCH_DATA, params)
            else:
                rows = []
            data = [{"partName": r[0], "month": r[1], "failureCount": r[2], "description": r[3]} for r in rows]
            if month and "All" not in month:
                data = [d for d in data if d['month'] in month]
            return data
        except Exception as e:
            logger.error(f"GNOVAC data lookup error: {e}")
            return {"error": str(e)}

    def get_gnovac_dashboard_data(self, user_id: int, part_name=None, month=None, base_model=None, mis_bucket=None, mfg_qtr=None) -> Dict[str, Any]:
        result = {"mfgMonth": [], "reportingMonth": [], "kms": [], "region": []}
        search_terms = [f"%{p.lower().replace(' ', '')}%" for p in part_name if p] if part_name else None
        params = {
            "user_id": user_id,
            "month_val": month if month and "All" not in month else None,
            "base_model": base_model if base_model and "All" not in base_model else None,
            "mis_bucket": mis_bucket if mis_bucket and "All" not in mis_bucket else None,
            "mfg_qtr": mfg_qtr if mfg_qtr and "All" not in mfg_qtr else None,
            "search_terms": search_terms,
        }
        try:
            range_row = self.db.execute_query(PartLabelerQueries.GNOVAC_GET_MFG_DATE_RANGE, {"user_id": user_id})
            if range_row and range_row[0][0] and range_row[0][1]:
                min_date, max_date = range_row[0]
                sequence = generate_month_sequence(min_date, max_date)
                db_data = {}
                if search_terms:
                    db_data = {r[0]: r[1] for r in self.db.execute_query(PartLabelerQueries.GNOVAC_GET_DASHBOARD_MFG_MONTH, params)}
                result["mfgMonth"] = [{"label": m, "value": db_data.get(m, 0)} for m in sequence]
        except Exception as e:
            logger.error(f"GNOVAC mfgMonth error: {e}")
        try:
            result["reportingMonth"] = [{"label": r[0], "value": r[1]} for r in self.db.execute_query(PartLabelerQueries.GNOVAC_GET_DASHBOARD_REPORTING_MONTH, params)]
        except Exception as e:
            logger.error(f"GNOVAC reportingMonth error: {e}")
        try:
            result["kms"] = [{"label": r[0], "value": r[1]} for r in self.db.execute_query(PartLabelerQueries.GNOVAC_GET_DASHBOARD_POINTER, params)]
        except Exception as e:
            logger.error(f"GNOVAC pointer error: {e}")
        try:
            result["region"] = [{"label": r[0], "value": r[1]} for r in self.db.execute_query(PartLabelerQueries.GNOVAC_GET_DASHBOARD_LOCATION, params)]
        except Exception as e:
            logger.error(f"GNOVAC location error: {e}")
        return result

    # =====================================================
    # RFI METHODS
    # =====================================================

    def process_mapped_rfi_data(self, file_path: str, mapping: Dict[str, str], user_id: int) -> int:
        try:
            df = pd.read_csv(file_path) if file_path.endswith('.csv') else pd.read_excel(file_path)
            all_cols = [
                'date_col', 'mfg_month', 'mfg_quarter', 'plant_name', 'vin_no', 'biw_no', 'model_name',
                'variant', 'fuel', 'drive_name', 'build_phase_name', 'software_v_name', 'color_name',
                'country_name', 'area_name', 'part_name', 'defect_name', 'location_name',
                'defect_type_name', 'severity_name', 'attribution_name', 'stage_name', 'root_cause',
                'ica', 'pca', 'target_date', 'responsibility', 'status', 'category_name', 'analysis_name',
                'action_plan_status', 'frequency'
            ]
            self.db.execute_update("DELETE FROM raw_rfi_data WHERE user_id = :user_id", {"user_id": user_id})
            records = []
            for _, row in df.iterrows():
                record = {col: '' for col in all_cols}
                record['user_id'] = user_id
                for db_col, user_col in mapping.items():
                    if user_col in row:
                        record[db_col] = safe_str(row[user_col])
                record['mfg_month'] = derive_mfg_month(record['date_col'])
                record['mfg_quarter'] = derive_mfg_quarter(record['date_col'])
                records.append(record)
                if len(records) >= 500:
                    self._insert_batch_rfi(records)
                    records = []
            if records:
                self._insert_batch_rfi(records)
            return len(df)
        except Exception as e:
            logger.error(f"Error processing RFI data: {e}")
            raise

    def _insert_batch_rfi(self, records):
        with self.db.get_session() as session:
            session.execute(text(PartLabelerQueries.INSERT_RAW_RFI), records)

    def get_rfi_filter_options(self, user_id: int) -> Dict[str, List[str]]:
        params = {"user_id": user_id}
        return {
            "models": [r[0] for r in self.db.execute_query(PartLabelerQueries.RFI_GET_UNIQUE_MODELS, params)],
            "mis_buckets": [r[0] for r in self.db.execute_query(PartLabelerQueries.RFI_GET_UNIQUE_MIS, params)],
            "mfg_quarters": [r[0] for r in self.db.execute_query(PartLabelerQueries.RFI_GET_UNIQUE_MFG_QUARTERS, params)],
            "mfg_months": [r[0] for r in self.db.execute_query(PartLabelerQueries.RFI_GET_UNIQUE_MFG_MONTHS, params)],
            "defect_types": [r[0] for r in self.db.execute_query(PartLabelerQueries.RFI_GET_UNIQUE_DEFECT_TYPES, params)],
        }

    def get_rfi_data(self, user_id: int, part_name=None, month=None, base_model=None, mis_bucket=None, mfg_qtr=None, defect_type=None) -> Any:
        try:
            params = {
                "base_model": base_model if base_model and "All" not in base_model else None,
                "mis_bucket": mis_bucket if mis_bucket and "All" not in mis_bucket else None,
                "mfg_qtr": mfg_qtr if mfg_qtr and "All" not in mfg_qtr else None,
                "defect_type": defect_type if defect_type and "All" not in defect_type else None,
                "user_id": user_id,
            }
            if part_name:
                params["search_term"] = f"%{part_name.lower().replace(' ', '')}%"
                rows = self.db.execute_query(PartLabelerQueries.RFI_SEARCH_DATA, params)
            else:
                rows = []
            data = [{"partName": r[0], "month": r[1], "failureCount": r[2], "description": r[3]} for r in rows]
            if month and "All" not in month:
                data = [d for d in data if d['month'] in month]
            return data
        except Exception as e:
            logger.error(f"RFI data lookup error: {e}")
            return {"error": str(e)}

    def get_rfi_dashboard_data(self, user_id: int, part_name=None, month=None, base_model=None, mis_bucket=None, mfg_qtr=None, defect_type=None) -> Dict[str, Any]:
        result = {"mfgMonth": [], "reportingMonth": [], "kms": [], "region": []}
        search_terms = [f"%{p.lower().replace(' ', '')}%" for p in part_name if p] if part_name else None
        params = {
            "user_id": user_id,
            "month_val": month if month and "All" not in month else None,
            "base_model": base_model if base_model and "All" not in base_model else None,
            "mis_bucket": mis_bucket if mis_bucket and "All" not in mis_bucket else None,
            "mfg_qtr": mfg_qtr if mfg_qtr and "All" not in mfg_qtr else None,
            "defect_type": defect_type if defect_type and "All" not in defect_type else None,
            "search_terms": search_terms,
        }
        try:
            range_row = self.db.execute_query(PartLabelerQueries.RFI_GET_MFG_DATE_RANGE, {"user_id": user_id})
            if range_row and range_row[0][0] and range_row[0][1]:
                min_date, max_date = range_row[0]
                sequence = generate_month_sequence(min_date, max_date)
                db_data = {}
                if search_terms:
                    db_data = {r[0]: r[1] for r in self.db.execute_query(PartLabelerQueries.RFI_GET_DASHBOARD_MFG_MONTH, params)}
                result["mfgMonth"] = [{"label": m, "value": db_data.get(m, 0)} for m in sequence]
        except Exception as e:
            logger.error(f"RFI mfgMonth error: {e}")
        try:
            result["reportingMonth"] = [{"label": r[0], "value": r[1]} for r in self.db.execute_query(PartLabelerQueries.RFI_GET_DASHBOARD_REPORTING_MONTH, params)]
        except Exception as e:
            logger.error(f"RFI reportingMonth error: {e}")
        try:
            result["kms"] = [{"label": r[0], "value": r[1]} for r in self.db.execute_query(PartLabelerQueries.RFI_GET_DASHBOARD_SEVERITY_DEFECTTYPE, params)]
        except Exception as e:
            logger.error(f"RFI severity/defecttype error: {e}")
        try:
            result["region"] = [{"label": r[0], "value": r[1]} for r in self.db.execute_query(PartLabelerQueries.RFI_GET_DASHBOARD_LOCATION, params)]
        except Exception as e:
            logger.error(f"RFI location error: {e}")
        return result

    # =====================================================
    # e-SQA METHODS
    # =====================================================

    def process_mapped_esqa_data(self, file_path: str, mapping: Dict[str, str], user_id: int) -> int:
        try:
            df = pd.read_csv(file_path) if file_path.endswith('.csv') else pd.read_excel(file_path)
            all_cols = [
                'concern_report_date', 'mfg_month', 'mfg_quarter', 'concern_number', 'pu_name',
                'concern_source', 'part_no', 'part_name', 'vendor_code', 'vendor_name',
                'concern_description', 'vehicle_model', 'vehicle_variant', 'concern_repeat',
                'concern_category', 'concern_severity', 'qty_reported', 'commodity',
                'concern_attribution', 'initial_analysis', 'sqa_officer', 'ica_possible',
                'reason_ica_not_possible', 'ica_details', 'ica_failure', 'segregation_qty',
                'ok_qty', 'rejection_qty', 'scrap_qty', 'rework_qty', 'deviation_qty',
                'line_loss', 'yard_hold', 'esqa_entry_required', 'justification_esqa',
                'esqa_number', 'esqa_posting_date'
            ]
            self.db.execute_update("DELETE FROM raw_esqa_data WHERE user_id = :user_id", {"user_id": user_id})
            records = []
            for _, row in df.iterrows():
                record = {col: '' for col in all_cols}
                record['user_id'] = user_id
                for db_col, user_col in mapping.items():
                    if user_col in row:
                        record[db_col] = safe_str(row[user_col])
                record['mfg_month'] = derive_mfg_month(record['concern_report_date'])
                record['mfg_quarter'] = derive_mfg_quarter(record['concern_report_date'])
                records.append(record)
                if len(records) >= 500:
                    self._insert_batch_esqa(records)
                    records = []
            if records:
                self._insert_batch_esqa(records)
            return len(df)
        except Exception as e:
            logger.error(f"Error processing e-SQA data: {e}")
            raise

    def _insert_batch_esqa(self, records):
        with self.db.get_session() as session:
            session.execute(text(PartLabelerQueries.INSERT_RAW_ESQA), records)

    def get_esqa_filter_options(self, user_id: int) -> Dict[str, List[str]]:
        params = {"user_id": user_id}
        return {
            "models": [r[0] for r in self.db.execute_query(PartLabelerQueries.ESQA_GET_UNIQUE_MODELS, params)],
            "mis_buckets": [r[0] for r in self.db.execute_query(PartLabelerQueries.ESQA_GET_UNIQUE_MIS, params)],
            "mfg_quarters": [r[0] for r in self.db.execute_query(PartLabelerQueries.ESQA_GET_UNIQUE_MFG_QUARTERS, params)],
            "mfg_months": [r[0] for r in self.db.execute_query(PartLabelerQueries.ESQA_GET_UNIQUE_MFG_MONTHS, params)],
        }

    def get_esqa_data(self, user_id: int, part_name=None, month=None, base_model=None, mis_bucket=None, mfg_qtr=None) -> Any:
        try:
            params = {
                "base_model": base_model if base_model and "All" not in base_model else None,
                "mis_bucket": mis_bucket if mis_bucket and "All" not in mis_bucket else None,
                "mfg_qtr": mfg_qtr if mfg_qtr and "All" not in mfg_qtr else None,
                "user_id": user_id,
            }
            if part_name:
                params["search_term"] = f"%{part_name.lower().replace(' ', '')}%"
                rows = self.db.execute_query(PartLabelerQueries.ESQA_SEARCH_DATA, params)
            else:
                rows = []
            data = [{"partName": r[0], "month": r[1], "failureCount": r[2], "description": r[3]} for r in rows]
            if month and "All" not in month:
                data = [d for d in data if d['month'] in month]
            return data
        except Exception as e:
            logger.error(f"e-SQA data lookup error: {e}")
            return {"error": str(e)}

    def get_esqa_dashboard_data(self, user_id: int, part_name=None, month=None, base_model=None, mis_bucket=None, mfg_qtr=None) -> Dict[str, Any]:
        result = {"mfgMonth": [], "reportingMonth": [], "kms": [], "region": []}
        search_terms = [f"%{p.lower().replace(' ', '')}%" for p in part_name if p] if part_name else None
        params = {
            "user_id": user_id,
            "month_val": month if month and "All" not in month else None,
            "base_model": base_model if base_model and "All" not in base_model else None,
            "mis_bucket": mis_bucket if mis_bucket and "All" not in mis_bucket else None,
            "mfg_qtr": mfg_qtr if mfg_qtr and "All" not in mfg_qtr else None,
            "search_terms": search_terms,
        }
        try:
            range_row = self.db.execute_query(PartLabelerQueries.ESQA_GET_MFG_DATE_RANGE, {"user_id": user_id})
            if range_row and range_row[0][0] and range_row[0][1]:
                min_date, max_date = range_row[0]
                sequence = generate_month_sequence(min_date, max_date)
                db_data = {}
                if search_terms:
                    db_data = {r[0]: r[1] for r in self.db.execute_query(PartLabelerQueries.ESQA_GET_DASHBOARD_MFG_MONTH, params)}
                result["mfgMonth"] = [{"label": m, "value": db_data.get(m, 0)} for m in sequence]
        except Exception as e:
            logger.error(f"e-SQA mfgMonth error: {e}")
        try:
            result["reportingMonth"] = [{"label": r[0], "value": r[1]} for r in self.db.execute_query(PartLabelerQueries.ESQA_GET_DASHBOARD_REPORTING_MONTH, params)]
        except Exception as e:
            logger.error(f"e-SQA reportingMonth error: {e}")
        try:
            result["kms"] = [{"label": r[0], "value": r[1]} for r in self.db.execute_query(PartLabelerQueries.ESQA_GET_DASHBOARD_CONCERN_SOURCE, params)]
        except Exception as e:
            logger.error(f"e-SQA concern source error: {e}")
        try:
            result["region"] = [{"label": r[0], "value": r[1]} for r in self.db.execute_query(PartLabelerQueries.ESQA_GET_DASHBOARD_LOCATION, params)]
        except Exception as e:
            logger.error(f"e-SQA location error: {e}")
        return result

    def get_dashboard_data(self, user_id: int, part_name: Optional[list[str]] = None, month: Optional[list[str]] = None, base_model: Optional[list[str]] = None, mis_bucket: Optional[list[str]] = None, mfg_qtr: Optional[list[str]] = None) -> Dict[str, Any]:
        """Fetch aggregated data for the 4 dashboard charts"""
        result = {"mfgMonth": [], "reportingMonth": [], "kms": [], "region": []}
        try:
            # Prepare search terms for multiple parts
            search_terms = None
            if part_name:
                search_terms = [f"%{p.lower().replace(' ', '')}%" for p in part_name if p]

            # Shared parameters
            params = {
                "user_id": user_id,
                "month_val": month if month and "All" not in month else None,
                "base_model": base_model if base_model and "All" not in base_model else None,
                "mis_bucket": mis_bucket if mis_bucket and "All" not in mis_bucket else None,
                "mfg_qtr": mfg_qtr if mfg_qtr and "All" not in mfg_qtr else None,
                "search_terms": search_terms
            }

            # 1. MFG Month
            try:
                range_row = self.db.execute_query(PartLabelerQueries.GET_MFG_DATE_RANGE, {"user_id": user_id})
                if range_row and range_row[0][0] and range_row[0][1]:
                    min_date, max_date = range_row[0]
                    sequence = generate_month_sequence(min_date, max_date)
                    db_data = {}
                    if search_terms:
                        raw = self.db.execute_query(PartLabelerQueries.GET_DASHBOARD_MFG_MONTH, params)
                        db_data = {normalize_month_label(r[0]): r[1] for r in raw}
                    result["mfgMonth"] = [{"label": m, "value": db_data.get(m, 0)} for m in sequence]
                else:
                    result["mfgMonth"] = []
            except Exception as e:
                logger.error(f"Error fetching MFG Month dashboard data: {e}")

            # 2. Reporting Month
            try:
                range_row_rep = self.db.execute_query(PartLabelerQueries.GET_REPORTING_DATE_RANGE, {"user_id": user_id})
                if range_row_rep and range_row_rep[0][0] and range_row_rep[0][1]:
                    min_date, max_date = range_row_rep[0]
                    sequence = generate_month_sequence(min_date, max_date)
                    db_data = {}
                    if search_terms:
                        raw = self.db.execute_query(PartLabelerQueries.GET_DASHBOARD_REPORTING_MONTH, params)
                        db_data = {normalize_month_label(r[0]): r[1] for r in raw}
                    result["reportingMonth"] = [{"label": m, "value": db_data.get(m, 0)} for m in sequence]
                else:
                    result["reportingMonth"] = []
            except Exception as e:
                logger.error(f"Error fetching Reporting Month dashboard data: {e}")

            # 3. KMS
            try:
                fixed_buckets = ['0-1k', '1k-2k', '2k-5k', '5k-10k', '10k-20k', '20k-30k', '30k above']
                db_data = {}
                if search_terms:
                    db_data = {r[0]: r[1] for r in self.db.execute_query(PartLabelerQueries.GET_DASHBOARD_KMS, params)}
                result["kms"] = [{"label": b, "value": db_data.get(b, 0)} for b in fixed_buckets]
            except Exception as e:
                logger.error(f"Error fetching KMS dashboard data: {e}")

            # 4. Region
            try:
                if search_terms:
                    result["region"] = [{"label": r[0], "value": r[1]} for r in self.db.execute_query(PartLabelerQueries.GET_DASHBOARD_REGION, params)]
                else:
                    regions = self.db.execute_query("SELECT DISTINCT region FROM raw_warranty_data WHERE user_id = :user_id AND region IS NOT NULL", {"user_id": user_id})
                    result["region"] = [{"label": r[0], "value": 0} for r in regions]
            except Exception as e:
                logger.error(f"Error fetching Region dashboard data: {e}")

            return result
        except Exception as e:
            logger.error(f"Global error fetching dashboard data: {e}")
            return result

    # =====================================================
    # EV WARRANTY METHODS
    # =====================================================

    def process_mapped_ev_data(self, file_path: str, mapping: Dict[str, str], user_id: int) -> int:
        try:
            df = pd.read_csv(file_path) if file_path.endswith('.csv') else pd.read_excel(file_path)
            all_cols = [
                'report_date', 'mfg_month', 'mfg_quarter', 'reporting_month', 'dealer_name',
                'location', 'state', 'brc_location', 'tar_no', 'vin_number', 'model',
                'problem_description', 'tekline_name', 'failure_category', 'rca', 'remarks_brc',
                'warranty_insurance', 'reason_of_failure', 'battery_motor', 'pvt_id', 'pvt_sno',
                'cause_pvt_list', 'kms', 'km_range', 'repaired_replaced', 'quality', 'repair_status',
                'part_updated', 'retail_date', 'no_of_days', 'mis', 'failure_mode',
            ]
            records = []
            for _, row in df.iterrows():
                record = {col: '' for col in all_cols}
                record['user_id'] = user_id
                for db_col, user_col in mapping.items():
                    if user_col in row:
                        record[db_col] = safe_str(row[user_col]).strip()
                record['mfg_month'] = derive_mfg_month(record['mfg_month'] or record['report_date'])
                record['mfg_quarter'] = derive_mfg_quarter(record['mfg_month'])
                record['reporting_month'] = derive_mfg_month(record['reporting_month'])
                record['model'] = record['model'].strip()
                records.append(record)
                if len(records) >= 500:
                    self._insert_batch_ev(records)
                    records = []
            if records:
                self._insert_batch_ev(records)
            return len(df)
        except Exception as e:
            logger.error(f"Error processing EV data: {e}")
            raise

    def _insert_batch_ev(self, records):
        with self.db.get_session() as session:
            session.execute(text(PartLabelerQueries.INSERT_RAW_EV), records)

    def get_ev_filter_options(self, user_id: int) -> Dict[str, List[str]]:
        params = {"user_id": user_id}
        return {
            "models": [r[0] for r in self.db.execute_query(PartLabelerQueries.EV_GET_UNIQUE_MODELS, params) if r[0]],
            "mis_buckets": [r[0] for r in self.db.execute_query(PartLabelerQueries.EV_GET_UNIQUE_MIS, params)],
            "mfg_quarters": [r[0] for r in self.db.execute_query(PartLabelerQueries.EV_GET_UNIQUE_MFG_QUARTERS, params)],
            "mfg_months": [r[0] for r in self.db.execute_query(PartLabelerQueries.EV_GET_UNIQUE_MFG_MONTHS, params)],
            "battery_motor_options": [r[0] for r in self.db.execute_query(PartLabelerQueries.EV_GET_UNIQUE_BATTERY_MOTOR, params)],
        }

    def get_ev_data(self, user_id: int, part_name=None, month=None, base_model=None, mis_bucket=None, mfg_qtr=None, battery_motor=None) -> Any:
        try:
            params = {
                "base_model": base_model if base_model and "All" not in base_model else None,
                "mis_bucket": mis_bucket if mis_bucket and "All" not in mis_bucket else None,
                "mfg_qtr": mfg_qtr if mfg_qtr and "All" not in mfg_qtr else None,
                "battery_motor": battery_motor if battery_motor and "All" not in battery_motor else None,
                "user_id": user_id,
            }
            if part_name:
                params["search_term"] = f"%{part_name.lower().replace(' ', '')}%"
                rows = self.db.execute_query(PartLabelerQueries.EV_SEARCH_DATA, params)
            else:
                rows = []
            data = [{"partName": r[0], "month": r[1], "failureCount": r[2], "description": r[3]} for r in rows]
            if month and "All" not in month:
                data = [d for d in data if d['month'] in month]
            return data
        except Exception as e:
            logger.error(f"EV data lookup error: {e}")
            return {"error": str(e)}

    def get_ev_dashboard_data(self, user_id: int, part_name=None, month=None, base_model=None, mis_bucket=None, mfg_qtr=None, battery_motor=None) -> Dict[str, Any]:
        result = {"mfgMonth": [], "reportingMonth": [], "kms": [], "region": []}
        search_terms = [f"%{p.lower().replace(' ', '')}%" for p in part_name if p] if part_name else None
        params = {
            "user_id": user_id,
            "month_val": month if month and "All" not in month else None,
            "base_model": base_model if base_model and "All" not in base_model else None,
            "mis_bucket": mis_bucket if mis_bucket and "All" not in mis_bucket else None,
            "mfg_qtr": mfg_qtr if mfg_qtr and "All" not in mfg_qtr else None,
            "battery_motor": battery_motor if battery_motor and "All" not in battery_motor else None,
            "search_terms": search_terms,
        }
        try:
            range_row = self.db.execute_query(
                "SELECT MIN(TO_DATE(mfg_month, 'Mon-YY')), MAX(TO_DATE(mfg_month, 'Mon-YY')) FROM raw_ev_data WHERE user_id = :user_id AND mfg_month IS NOT NULL AND mfg_month != ''",
                {"user_id": user_id}
            )
            if range_row and range_row[0][0] and range_row[0][1]:
                min_date, max_date = range_row[0]
                sequence = generate_month_sequence(min_date, max_date)
                db_data = {}
                if search_terms:
                    db_data = {r[0]: r[1] for r in self.db.execute_query(PartLabelerQueries.EV_GET_DASHBOARD_MFG_MONTH, params)}
                result["mfgMonth"] = [{"label": m, "value": db_data.get(m, 0)} for m in sequence]
        except Exception as e:
            logger.error(f"EV mfgMonth error: {e}")
        try:
            result["reportingMonth"] = [{"label": r[0], "value": r[1]} for r in self.db.execute_query(PartLabelerQueries.EV_GET_DASHBOARD_REPORTING_MONTH, params)]
        except Exception as e:
            logger.error(f"EV reportingMonth error: {e}")
        try:
            result["kms"] = [{"label": r[0], "value": r[1]} for r in self.db.execute_query(PartLabelerQueries.EV_GET_DASHBOARD_KM_RANGE, params)]
        except Exception as e:
            logger.error(f"EV km range error: {e}")
        try:
            result["region"] = [{"label": r[0], "value": r[1]} for r in self.db.execute_query(PartLabelerQueries.EV_GET_DASHBOARD_STATE, params)]
        except Exception as e:
            logger.error(f"EV state error: {e}")
        return result
